"""Shared fixtures: python-pptx builds fixture decks, pptx_rs exercises them."""

from __future__ import annotations

from pathlib import Path

import pytest
from pptx import Presentation as PptxPresentation
from pptx.util import Inches


@pytest.fixture
def fixture_deck(tmp_path: Path) -> Path:
    """A 3-slide deck built with python-pptx (the reference implementation)."""
    prs = PptxPresentation()
    blank = prs.slide_layouts[6]
    for i in range(3):
        slide = prs.slides.add_slide(blank)
        for j in range(2):
            box = slide.shapes.add_textbox(Inches(0.5 + j), Inches(0.5), Inches(2), Inches(1))
            tf = box.text_frame
            tf.text = f"slide {i} box {j}"
            para = tf.add_paragraph()
            para.add_run().text = f"para 2 of box {j} & <special> chars"
    path = tmp_path / "fixture.pptx"
    prs.save(str(path))
    return path
