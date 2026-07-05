"""Reading decks produced by python-pptx."""

from __future__ import annotations

import io
from pathlib import Path

import pytest

import pptx_rs


def test_opens_default_template_when_no_path_given():
    prs = pptx_rs.Presentation()
    assert len(prs.slides) == 0
    assert prs.slide_width == 9144000
    assert prs.slide_height == 6858000


def test_opens_from_path_str_and_pathlib(fixture_deck: Path):
    assert len(pptx_rs.Presentation(str(fixture_deck)).slides) == 3
    assert len(pptx_rs.Presentation(fixture_deck).slides) == 3


def test_opens_from_file_like(fixture_deck: Path):
    with open(fixture_deck, "rb") as f:
        prs = pptx_rs.Presentation(f)
    assert len(prs.slides) == 3


def test_slides_are_iterable_and_indexable(fixture_deck: Path):
    prs = pptx_rs.Presentation(fixture_deck)
    assert len(list(prs.slides)) == 3
    assert prs.slides[0].slide_id == 256
    assert prs.slides[-1].slide_id == prs.slides[2].slide_id
    with pytest.raises(IndexError):
        prs.slides[3]


def test_shapes_expose_identity_and_geometry(fixture_deck: Path):
    shape = pptx_rs.Presentation(fixture_deck).slides[0].shapes[0]
    assert shape.shape_id == 2
    assert shape.name == "TextBox 1"
    assert shape.has_text_frame
    assert shape.left == 457200  # Inches(0.5)
    assert shape.top == 457200
    assert shape.width == 1828800
    assert shape.height == 914400


def test_text_traversal_matches_python_pptx(fixture_deck: Path):
    from pptx import Presentation as PptxPresentation

    ours = pptx_rs.Presentation(fixture_deck)
    theirs = PptxPresentation(str(fixture_deck))
    for our_slide, their_slide in zip(ours.slides, theirs.slides, strict=True):
        for our_shape, their_shape in zip(our_slide.shapes, their_slide.shapes, strict=True):
            assert our_shape.has_text_frame == their_shape.has_text_frame
            if not our_shape.has_text_frame:
                continue
            our_tf, their_tf = our_shape.text_frame, their_shape.text_frame
            assert our_tf.text == their_tf.text
            for our_p, their_p in zip(our_tf.paragraphs, their_tf.paragraphs, strict=True):
                assert our_p.text == their_p.text
                our_runs = [r.text for r in our_p.runs]
                their_runs = [r.text for r in their_p.runs]
                assert our_runs == their_runs


def test_escaped_characters_survive(fixture_deck: Path):
    para = pptx_rs.Presentation(fixture_deck).slides[0].shapes[0].text_frame.paragraphs[1]
    assert para.text == "para 2 of box 0 & <special> chars"


def test_save_to_file_like(fixture_deck: Path):
    prs = pptx_rs.Presentation(fixture_deck)
    buf = io.BytesIO()
    prs.save(buf)
    assert buf.getvalue()[:2] == b"PK"
