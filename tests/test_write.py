"""Editing and writing: every result is re-validated with python-pptx."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation as PptxPresentation

import pptx_rs
from pptx_rs.util import Emu, Inches


def reopen_with_python_pptx(prs: pptx_rs._core.Presentation, path: Path):
    prs.save(path)
    return PptxPresentation(str(path))


def test_run_text_edit_round_trips(fixture_deck: Path, tmp_path: Path):
    prs = pptx_rs.Presentation(fixture_deck)
    prs.slides[0].shapes[0].text_frame.paragraphs[0].runs[0].text = "edited & <checked>"

    verified = reopen_with_python_pptx(prs, tmp_path / "out.pptx")
    text = verified.slides[0].shapes[0].text_frame.paragraphs[0].runs[0].text
    assert text == "edited & <checked>"


def test_text_frame_text_setter_creates_paragraphs(fixture_deck: Path, tmp_path: Path):
    prs = pptx_rs.Presentation(fixture_deck)
    prs.slides[0].shapes[0].text_frame.text = "line1\nline2\vsoft-break"

    verified = reopen_with_python_pptx(prs, tmp_path / "out.pptx")
    tf = verified.slides[0].shapes[0].text_frame
    assert tf.text == "line1\nline2\vsoft-break"
    assert len(tf.paragraphs) == 2


def test_add_textbox_matches_python_pptx_semantics(fixture_deck: Path, tmp_path: Path):
    prs = pptx_rs.Presentation(fixture_deck)
    shapes = prs.slides[0].shapes
    box = shapes.add_textbox(Inches(1), Inches(2), Inches(3), Inches(0.5))
    box.text_frame.text = "new box"

    verified = reopen_with_python_pptx(prs, tmp_path / "out.pptx")
    added = verified.slides[0].shapes[-1]
    assert added.shape_id == 4  # ids 2..3 already taken by fixture boxes
    assert added.name == "TextBox 3"
    assert added.left == Inches(1)
    assert added.top == Inches(2)
    assert added.text_frame.text == "new box"


def test_shape_ids_stay_unique_across_many_adds(fixture_deck: Path, tmp_path: Path):
    prs = pptx_rs.Presentation(fixture_deck)
    shapes = prs.slides[1].shapes
    for _ in range(50):
        shapes.add_textbox(Emu(0), Emu(0), Emu(914400), Emu(914400))

    verified = reopen_with_python_pptx(prs, tmp_path / "out.pptx")
    ids = [s.shape_id for s in verified.slides[1].shapes]
    assert len(ids) == len(set(ids))


def test_geometry_setters(fixture_deck: Path, tmp_path: Path):
    prs = pptx_rs.Presentation(fixture_deck)
    shape = prs.slides[0].shapes[0]
    shape.left = Inches(2)
    shape.width = Inches(4)

    verified = reopen_with_python_pptx(prs, tmp_path / "out.pptx")
    assert verified.slides[0].shapes[0].left == Inches(2)
    assert verified.slides[0].shapes[0].width == Inches(4)


def test_add_paragraph_and_run(fixture_deck: Path, tmp_path: Path):
    prs = pptx_rs.Presentation(fixture_deck)
    tf = prs.slides[0].shapes[0].text_frame
    para = tf.add_paragraph()
    run = para.add_run()
    run.text = "appended"

    verified = reopen_with_python_pptx(prs, tmp_path / "out.pptx")
    assert verified.slides[0].shapes[0].text_frame.paragraphs[-1].text == "appended"


def test_untouched_parts_round_trip_byte_exact(fixture_deck: Path, tmp_path: Path):
    import zipfile

    prs = pptx_rs.Presentation(fixture_deck)
    prs.slides[0].shapes[0].text_frame.text = "touch only slide 1"
    out = tmp_path / "out.pptx"
    prs.save(out)

    with zipfile.ZipFile(fixture_deck) as src, zipfile.ZipFile(out) as dst:
        assert src.namelist() == dst.namelist()
        for name in src.namelist():
            if name == "ppt/slides/slide1.xml":
                continue
            assert src.read(name) == dst.read(name), f"{name} changed unexpectedly"
