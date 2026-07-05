"""add_slide + slide_layouts: every result re-validated with python-pptx."""

from __future__ import annotations

from pathlib import Path

import pytest
from pptx import Presentation as PptxPresentation

import pptx_rs
from pptx_rs.util import Inches


def test_slide_layouts_match_python_pptx():
    ours = pptx_rs.Presentation()
    theirs = PptxPresentation()
    our_names = [layout.name for layout in ours.slide_layouts]
    their_names = [layout.name for layout in theirs.slide_layouts]
    assert our_names == their_names
    assert len(ours.slide_masters) == 1


def test_get_by_name_and_index():
    prs = pptx_rs.Presentation()
    blank = prs.slide_layouts.get_by_name("Blank")
    assert blank is not None
    assert prs.slide_layouts.index(blank) == 6
    assert prs.slide_layouts.get_by_name("No Such Layout") is None


def test_add_slide_to_default_template(tmp_path: Path):
    prs = pptx_rs.Presentation()
    layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(layout)
    assert len(prs.slides) == 1
    assert slide.slide_id == 256
    assert slide.slide_layout == layout

    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
    box.text_frame.text = "from scratch"
    out = tmp_path / "out.pptx"
    prs.save(out)

    verified = PptxPresentation(str(out))
    assert len(verified.slides) == 1
    assert verified.slides[0].slide_layout.name == "Blank"
    assert verified.slides[0].shapes[0].text_frame.text == "from scratch"


def test_add_slide_clones_layout_placeholders(tmp_path: Path):
    prs = pptx_rs.Presentation()
    title_layout = prs.slide_layouts[0]  # Title Slide: ctrTitle + subTitle
    slide = prs.slides.add_slide(title_layout)
    slide.shapes[0].text_frame.text = "My Title"
    out = tmp_path / "out.pptx"
    prs.save(out)

    verified = PptxPresentation(str(out))
    placeholders = list(verified.slides[0].placeholders)
    assert len(placeholders) == 2
    assert verified.slides[0].shapes.title.text == "My Title"
    # latent placeholders (date/footer/slide-number) must not be cloned
    reference = PptxPresentation()
    ref_slide = reference.slides.add_slide(reference.slide_layouts[0])
    assert len(placeholders) == len(list(ref_slide.placeholders))


def test_add_many_slides_round_trips(tmp_path: Path, fixture_deck: Path):
    prs = pptx_rs.Presentation(fixture_deck)
    layout = prs.slide_layouts[6]
    for i in range(20):
        slide = prs.slides.add_slide(layout)
        slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1)).text_frame.text = f"added {i}"
    out = tmp_path / "out.pptx"
    prs.save(out)

    verified = PptxPresentation(str(out))
    assert len(verified.slides) == 23
    ids = [s.slide_id for s in verified.slides]
    assert len(ids) == len(set(ids))
    assert verified.slides[-1].shapes[0].text_frame.text == "added 19"


def test_added_slide_is_immediately_usable():
    prs = pptx_rs.Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    assert slide.shapes[0].has_text_frame
    assert prs.slides[-1] == slide


def test_layout_index_of_foreign_layout_raises():
    prs1 = pptx_rs.Presentation()
    prs2 = pptx_rs.Presentation()
    layout2 = prs2.slide_layouts[0]
    with pytest.raises(ValueError, match="not in this"):
        prs1.slide_layouts.index(layout2)
