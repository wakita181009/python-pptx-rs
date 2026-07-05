"""MarkItDown integration: pptx_rs must be a drop-in engine for python-pptx.

MarkItDown's PptxConverter resolves ``pptx.Presentation`` at call time, so
monkeypatching that one attribute swaps the whole parsing engine while
``MSO_SHAPE_TYPE`` comparisons keep using real python-pptx enums — exactly the
cross-library situation our IntEnum-based enums must survive.
"""

from __future__ import annotations

import io
from pathlib import Path

import pptx
import pytest
from markitdown import MarkItDown
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

import pptx_rs

MARKITDOWN_DECK = Path(__file__).parents[1] / "fixtures" / "markitdown" / "test.pptx"


def _convert(source: Path | io.BytesIO) -> str:
    md = MarkItDown()
    if isinstance(source, Path):
        with source.open("rb") as f:
            return md.convert_stream(f, file_extension=".pptx").markdown
    source.seek(0)
    return md.convert_stream(source, file_extension=".pptx").markdown


@pytest.fixture
def synthetic_deck() -> io.BytesIO:
    """A deck built with python-pptx covering title/table/chart/picture/notes."""
    prs = pptx.Presentation()

    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text_frame.text = "Synthetic Deck"
    title_slide.notes_slide.notes_text_frame.text = "speaker note 1\nspeaker note 2"

    table_slide = prs.slides.add_slide(prs.slide_layouts[6])
    table = table_slide.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(4), Inches(2)).table
    for r, row in enumerate(table.rows):
        for c, cell in enumerate(row.cells):
            cell.text = f"cell {r}{c}"

    chart_slide = prs.slides.add_slide(prs.slide_layouts[6])
    chart_data = CategoryChartData()
    chart_data.categories = ["East", "West"]
    chart_data.add_series("Sales", (1.2, 3.4))
    chart_slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(1), Inches(6), Inches(4), chart_data)

    buf = io.BytesIO()
    prs.save(buf)
    return buf


def test_markitdown_real_deck_output_identical(monkeypatch: pytest.MonkeyPatch) -> None:
    baseline = _convert(MARKITDOWN_DECK)
    monkeypatch.setattr(pptx, "Presentation", pptx_rs.Presentation)
    assert _convert(MARKITDOWN_DECK) == baseline


def test_markitdown_real_deck_content(monkeypatch: pytest.MonkeyPatch) -> None:
    monkeypatch.setattr(pptx, "Presentation", pptx_rs.Presentation)
    markdown = _convert(MARKITDOWN_DECK)
    assert "<!-- Slide number: 6 -->" in markdown  # all six slides converted
    assert "![" in markdown  # pictures
    assert "### Chart" in markdown  # chart rendered as a table
    assert "| Category | Series 1 |" in markdown


def test_markitdown_synthetic_deck_output_identical(
    synthetic_deck: io.BytesIO, monkeypatch: pytest.MonkeyPatch
) -> None:
    baseline = _convert(synthetic_deck)
    assert "speaker note 1" in baseline  # notes path is actually exercised
    monkeypatch.setattr(pptx, "Presentation", pptx_rs.Presentation)
    assert _convert(synthetic_deck) == baseline


def test_placeholder_geometry_matches_python_pptx() -> None:
    """Title-slide placeholders carry no direct a:xfrm, so left/top/width/height
    are resolved by inheritance (slide -> layout -> master). Regression guard for
    the inheritance walk mis-identifying a master's own layouts as its parent."""
    prs = pptx.Presentation()
    prs.slides.add_slide(prs.slide_layouts[0])
    buf = io.BytesIO()
    prs.save(buf)

    ref = pptx.Presentation(io.BytesIO(buf.getvalue()))
    ours = pptx_rs.Presentation(io.BytesIO(buf.getvalue()))
    geom = lambda deck: [  # noqa: E731
        (s.left, s.top, s.width, s.height) for s in deck.slides[0].shapes
    ]
    assert geom(ours) == geom(ref)


def test_chart_title_without_rich_text_returns_empty(monkeypatch: pytest.MonkeyPatch) -> None:
    """A chart with a <c:title> but no c:tx/c:rich must yield an empty title
    string, not raise -- otherwise markitdown's broad except drops the whole
    chart from the output."""
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE

    prs = pptx.Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    chart_data = CategoryChartData()
    chart_data.categories = ["a", "b"]
    chart_data.add_series("S", (1.0, 2.0))
    frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(1), Inches(5), Inches(4), chart_data
    )
    frame.chart.has_title = True
    buf = io.BytesIO()
    prs.save(buf)

    chart = pptx_rs.Presentation(io.BytesIO(buf.getvalue())).slides[0].shapes[0].chart
    assert chart.has_title is True
    assert chart.chart_title.text_frame.text == ""

    monkeypatch.setattr(pptx, "Presentation", pptx_rs.Presentation)
    assert "### Chart" in _convert(io.BytesIO(buf.getvalue()))


def test_shape_type_equals_python_pptx_enum() -> None:
    from pptx.enum.shapes import MSO_SHAPE_TYPE as PPTX_MSO

    from pptx_rs.enum.shapes import MSO, MSO_SHAPE_TYPE

    assert MSO is MSO_SHAPE_TYPE
    assert list(MSO_SHAPE_TYPE.__members__) == list(PPTX_MSO.__members__)
    assert all(MSO_SHAPE_TYPE[name] == PPTX_MSO[name] for name in MSO_SHAPE_TYPE.__members__)

    prs = pptx_rs.Presentation(MARKITDOWN_DECK)
    kinds = {shape.shape_type for slide in prs.slides for shape in slide.shapes}
    assert {
        PPTX_MSO.PLACEHOLDER,
        PPTX_MSO.PICTURE,
        PPTX_MSO.TABLE,
        PPTX_MSO.CHART,
        PPTX_MSO.GROUP,
    } <= kinds
