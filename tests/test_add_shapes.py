"""add_picture / add_shape / add_table: re-validated with python-pptx.

Native-size and naming behavior is asserted as parity against real
python-pptx performing the same operation.
"""

from __future__ import annotations

import io
import zipfile
from pathlib import Path

import pytest
from pptx import Presentation as PptxPresentation
from pptx.enum.shapes import MSO_SHAPE as REAL_MSO_SHAPE
from pptx.enum.shapes import MSO_SHAPE_TYPE as REAL_MSO_SHAPE_TYPE

import pptx_rs
from pptx_rs.enum.shapes import MSO_SHAPE
from pptx_rs.util import Inches

TEST_FILES = Path(__file__).parent / "compat" / "features" / "steps" / "test_files"


def reopen_with_python_pptx(prs: pptx_rs._core.Presentation, path: Path):
    prs.save(path)
    return PptxPresentation(str(path))


def new_blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


# ---------------------------------------------------------------------------
# add_picture
# ---------------------------------------------------------------------------


@pytest.mark.parametrize(
    ("filename", "expected_ext"),
    [
        ("monty-truth.png", "png"),
        ("python-icon.jpeg", "jpg"),
        ("sonic.gif", "gif"),
        ("python.bmp", "bmp"),
        ("72-dpi.tiff", "tiff"),
        ("CVS_LOGO.WMF", "wmf"),
        ("pic.emf", "wmf"),
    ],
)
def test_add_picture_native_size_matches_python_pptx(filename: str, expected_ext: str, tmp_path: Path):
    image = str(TEST_FILES / filename)

    prs = pptx_rs.Presentation()
    shapes = new_blank_slide(prs).shapes
    pic = shapes.add_picture(image, Inches(1), Inches(1))

    real_prs = PptxPresentation()
    real_shapes = real_prs.slides.add_slide(real_prs.slide_layouts[6]).shapes
    real_pic = real_shapes.add_picture(image, Inches(1), Inches(1))

    assert (pic.width, pic.height) == (real_pic.width, real_pic.height)
    assert pic.name == real_pic.name
    # shape factory returns the python-pptx proxy class name
    assert type(pic).__name__ == "Picture"

    verified = reopen_with_python_pptx(prs, tmp_path / "out.pptx")
    added = verified.slides[0].shapes[-1]
    assert added.shape_type == REAL_MSO_SHAPE_TYPE.PICTURE
    assert added.image.ext == expected_ext
    assert added.image.blob == Path(image).read_bytes()


def test_add_picture_scales_to_preserve_aspect_ratio():
    image = str(TEST_FILES / "monty-truth.png")

    prs = pptx_rs.Presentation()
    shapes = new_blank_slide(prs).shapes
    by_width = shapes.add_picture(image, 0, 0, width=Inches(2))
    by_height = shapes.add_picture(image, 0, 0, height=Inches(2))
    stretched = shapes.add_picture(image, 0, 0, width=Inches(2), height=Inches(2))

    real_prs = PptxPresentation()
    real_shapes = real_prs.slides.add_slide(real_prs.slide_layouts[6]).shapes
    real_by_width = real_shapes.add_picture(image, 0, 0, width=Inches(2))

    assert (by_width.width, by_width.height) == (real_by_width.width, real_by_width.height)
    assert by_height.height == Inches(2)
    assert (stretched.width, stretched.height) == (Inches(2), Inches(2))


def test_add_picture_accepts_file_like_object(tmp_path: Path):
    blob = (TEST_FILES / "monty-truth.png").read_bytes()

    prs = pptx_rs.Presentation()
    shapes = new_blank_slide(prs).shapes
    stream = io.BytesIO(blob)
    stream.read()  # cursor at EOF: add_picture must rewind
    shapes.add_picture(stream, Inches(1), Inches(1))

    verified = reopen_with_python_pptx(prs, tmp_path / "out.pptx")
    assert verified.slides[0].shapes[-1].image.blob == blob


def test_add_picture_dedupes_identical_images(tmp_path: Path):
    image = str(TEST_FILES / "monty-truth.png")

    prs = pptx_rs.Presentation()
    slide1 = new_blank_slide(prs)
    slide2 = new_blank_slide(prs)
    slide1.shapes.add_picture(image, 0, 0)
    slide1.shapes.add_picture(image, Inches(1), 0)
    slide2.shapes.add_picture(image, 0, 0)

    out = tmp_path / "out.pptx"
    prs.save(out)
    with zipfile.ZipFile(out) as zf:
        media = [n for n in zf.namelist() if n.startswith("ppt/media/")]
        assert media == ["ppt/media/image1.png"]
        # both p:pic on the first slide share the deduped relationship
        slide1_xml = zf.read("ppt/slides/slide1.xml")
        assert slide1_xml.count(b"r:embed") == 2
        rels = zf.read("ppt/slides/_rels/slide1.xml.rels")
        assert rels.count(b"image1.png") == 1

    verified = PptxPresentation(str(out))
    for slide in verified.slides:
        for shape in slide.shapes:
            assert shape.image.ext == "png"


def test_add_picture_registers_default_content_type(tmp_path: Path):
    prs = pptx_rs.Presentation()
    shapes = new_blank_slide(prs).shapes
    shapes.add_picture(str(TEST_FILES / "monty-truth.png"), 0, 0)

    out = tmp_path / "out.pptx"
    prs.save(out)
    with zipfile.ZipFile(out) as zf:
        content_types = zf.read("[Content_Types].xml").decode()
    assert 'Extension="png"' in content_types
    assert content_types.index('Extension="png"') < content_types.index("<Override")


def test_add_picture_rejects_unsupported_format():
    prs = pptx_rs.Presentation()
    shapes = new_blank_slide(prs).shapes
    with pytest.raises(ValueError, match="unsupported image format"):
        shapes.add_picture(io.BytesIO(b"plainly not an image"), 0, 0)


# ---------------------------------------------------------------------------
# add_shape
# ---------------------------------------------------------------------------


def test_add_shape_matches_python_pptx_semantics(tmp_path: Path):
    prs = pptx_rs.Presentation()
    shapes = new_blank_slide(prs).shapes
    shape = shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(2), Inches(3), Inches(1))
    assert shape.name == "Rounded Rectangle 1"

    verified = reopen_with_python_pptx(prs, tmp_path / "out.pptx")
    added = verified.slides[0].shapes[-1]
    assert added.shape_type == REAL_MSO_SHAPE_TYPE.AUTO_SHAPE
    assert added.auto_shape_type == REAL_MSO_SHAPE.ROUNDED_RECTANGLE
    assert (added.left, added.top) == (Inches(1), Inches(2))
    assert (added.width, added.height) == (Inches(3), Inches(1))
    assert added.has_text_frame


def test_add_shape_accepts_real_python_pptx_enum(tmp_path: Path):
    prs = pptx_rs.Presentation()
    shapes = new_blank_slide(prs).shapes
    shape = shapes.add_shape(REAL_MSO_SHAPE.NO_SYMBOL, 0, 0, Inches(1), Inches(1))
    # basename contains XML-hostile quotes; must round-trip escaped
    assert shape.name == '"No" Symbol 1'

    verified = reopen_with_python_pptx(prs, tmp_path / "out.pptx")
    assert verified.slides[0].shapes[-1].auto_shape_type == REAL_MSO_SHAPE.NO_SYMBOL


def test_add_shape_rejects_unknown_id():
    prs = pptx_rs.Presentation()
    shapes = new_blank_slide(prs).shapes
    with pytest.raises(KeyError):
        shapes.add_shape(99999, 0, 0, 1, 1)


# ---------------------------------------------------------------------------
# add_table
# ---------------------------------------------------------------------------


def test_add_table_matches_python_pptx_semantics(tmp_path: Path):
    prs = pptx_rs.Presentation()
    shapes = new_blank_slide(prs).shapes
    # width not divisible by cols: last column absorbs the remainder
    frame = shapes.add_table(2, 3, Inches(1), Inches(1), 1000000, 600000)
    frame.table.cell(0, 0).text = "Header & Co"
    frame.table.cell(1, 2).text = "multi\nline"
    assert frame.name == "Table 1"
    assert type(frame).__name__ == "GraphicFrame"

    verified = reopen_with_python_pptx(prs, tmp_path / "out.pptx")
    added = verified.slides[0].shapes[-1]
    assert added.has_table
    table = added.table
    assert (len(table.rows), len(table.columns)) == (2, 3)
    assert [c.width for c in table.columns] == [333333, 333333, 333334]
    assert [r.height for r in table.rows] == [300000, 300000]
    assert table.cell(0, 0).text == "Header & Co"
    assert table.cell(1, 2).text == "multi\nline"
    assert table.first_row and table.horz_banding


def test_add_table_cell_read_back_through_pptx_rs():
    prs = pptx_rs.Presentation()
    shapes = new_blank_slide(prs).shapes
    frame = shapes.add_table(2, 2, 0, 0, Inches(4), Inches(1))
    frame.table.cell(0, 1).text = "roundtrip"
    assert frame.table.cell(0, 1).text == "roundtrip"
    assert frame.table.rows[0].cells[1].text == "roundtrip"
    with pytest.raises(IndexError):
        frame.table.cell(2, 0)
