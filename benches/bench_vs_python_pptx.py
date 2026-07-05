"""Head-to-head benchmark: pptx_rs vs python-pptx on identical workloads."""

from __future__ import annotations

import statistics
import sys
import time

from pptx import Presentation as PyPresentation
from pptx.util import Inches as PyInches

import pptx_rs
from pptx_rs.util import Inches

N_SLIDES = 1000
SHAPES = 10
DECK = "/tmp/bench_deck.pptx"


def build_fixture() -> None:
    prs = PyPresentation()
    blank = prs.slide_layouts[6]
    for i in range(N_SLIDES):
        slide = prs.slides.add_slide(blank)
        for j in range(SHAPES):
            box = slide.shapes.add_textbox(PyInches(0.2), PyInches(0.2), PyInches(2), PyInches(1))
            box.text_frame.paragraphs[0].add_run().text = f"slide {i} shape {j} content text"
    prs.save(DECK)


def timeit(fn, repeat: int = 3) -> float:
    times = [0.0] * repeat
    for i in range(repeat):
        t0 = time.perf_counter()
        fn()
        times[i] = time.perf_counter() - t0
    return statistics.median(times)


def traverse(prs) -> int:
    n = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    n += len(run.text)
    return n


def bench_open(mod):
    return timeit(lambda: mod(DECK))


def bench_traverse(mod):
    prs = mod(DECK)
    return timeit(lambda: traverse(prs))


def bench_open_traverse(mod):
    return timeit(lambda: traverse(mod(DECK)))


def bench_edit_save(mod):
    def job():
        prs = mod(DECK)
        prs.slides[0].shapes[0].text_frame.paragraphs[0].runs[0].text = "edited"
        prs.save("/tmp/bench_out.pptx")

    return timeit(job)


def bench_add_shapes(mod, inches):
    def job():
        prs = mod(DECK)
        shapes = prs.slides[0].shapes
        for _ in range(500):
            shapes.add_textbox(inches(0.1), inches(0.1), inches(1), inches(0.5))

    return timeit(job)


def py_open(path):
    return PyPresentation(path)


def rs_open(path):
    return pptx_rs.Presentation(path)


def main() -> None:
    build_fixture()
    print(f"deck: {N_SLIDES} slides x {SHAPES} textboxes ({N_SLIDES * SHAPES} runs)\n")
    rows = [
        ("open", bench_open(py_open), bench_open(rs_open)),
        ("open + traverse all text", bench_open_traverse(py_open), bench_open_traverse(rs_open)),
        ("traverse (pre-opened)", bench_traverse(py_open), bench_traverse(rs_open)),
        ("open + 1 edit + save", bench_edit_save(py_open), bench_edit_save(rs_open)),
        (
            "add 500 textboxes to one slide",
            bench_add_shapes(py_open, PyInches),
            bench_add_shapes(rs_open, Inches),
        ),
    ]
    print(f"{'workload':<32} {'python-pptx':>12} {'pptx_rs':>10} {'speedup':>8}")
    for name, py_t, rs_t in rows:
        print(f"{name:<32} {py_t * 1000:>10.1f}ms {rs_t * 1000:>8.1f}ms {py_t / rs_t:>7.1f}x")


if __name__ == "__main__":
    sys.exit(main())
